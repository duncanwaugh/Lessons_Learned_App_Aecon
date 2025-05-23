import os
import hashlib
import streamlit as st
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from openai import OpenAI
from docxtpl import DocxTemplate, InlineImage
from docx import Document
from docx.shared import Inches
from PIL import Image as PILImage

# ─── Load API keys ────────────────────────────────────────────────────────────
load_dotenv()
OPENAI_KEY = os.getenv("OPENAI_API_KEY")
DEEPL_KEY  = os.getenv("DEEPL_API_KEY")
if not OPENAI_KEY:
    st.error("Missing OPENAI_API_KEY in .env or Streamlit secrets")
client = OpenAI(api_key=OPENAI_KEY)

# ─── Extract + dedupe text & images from PPTX ─────────────────────────────────
def extract_text_and_images_from_pptx(path: str):
    prs = Presentation(path)
    text, images, seen_hashes = "", [], set()
    os.makedirs("images", exist_ok=True)

    for idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if shp.has_text_frame:
                text += shp.text_frame.text + "\n"
            elif shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
            elif shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shp.image.blob
                h = hashlib.sha256(blob).hexdigest()
                if h in seen_hashes:
                    continue
                seen_hashes.add(h)
                ext = shp.image.ext
                fn  = f"images/slide{idx+1}_{len(images)}.{ext}"
                with open(fn, "wb") as imgf:
                    imgf.write(blob)
                images.append(fn)

    return text, images

# ─── OpenAI extraction with detailed summary ──────────────────────────────────
def summarize_and_extract(text: str) -> str:
    system_msg = ("You are a concise safety-report writer for Aecon."
                  "You use the presentation text given to write sophisticated")
    prompt = f"""
You are preparing a Lessons Learned report from a serious incident.  Use **plain**, **easy‑to‑understand** language and a **neutral, positive** tone.

for contributing factors and lessons learned use bullet lists only, one idea per line, avoid technical jargon and blaming, keep it factual, clear and solution oriented.
Produce each section clearly labeled with exactly these labels (and nothing else):

Title: (Keep it very concise (3–5 words) summarizing the core event)  
Aecon Business Sector:  
Project/Location:  
Date of Event:  
Event Type:  
Event Summary Header: (One detailed sentence that captures the essence of what happened) 
Event Summary:  (Write a detailed, multi‑paragraph narrative covering:  Background and context, Step‑by‑step sequence of events, Immediate outcome and injuries/damages, Broader impacts (delays, reputation, etc.))
Contributing Factors:  (avoid any technical terms)
Lessons Learned:   (avoid any technical terms)


Here is the presentation text to parse into those sections:
{text}
"""
    r = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role":"system","content":system_msg},
            {"role":"user","content":prompt}
        ],
        temperature=0.2,
        max_tokens=1500,
    )
    return r.choices[0].message.content.strip()

# ─── Translation helpers ──────────────────────────────────────────────────────
def translate_to_french_openai(text: str) -> str:
    prompt = f"Translate into professional French Canadian, keep formatting:\n\n{text}"
    r = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role":"user","content":prompt}],
        temperature=0.2,
        max_tokens=1000,
    )
    return r.choices[0].message.content.strip()


def translate_to_spanish_openai(text: str) -> str:
    prompt = f"Translate into professional Spanish, keep formatting, section headers and bullets intact:\n\n{text}"
    r = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role":"user","content":prompt}],
        temperature=0.2,
        max_tokens=1000,
    )
    return r.choices[0].message.content.strip()


def translate_to_deepl(text: str, lang: str) -> str:
    if not DEEPL_KEY:
        st.error("Missing DEEPL_API_KEY in .env or Streamlit secrets")
        return text
    target = "FR" if lang == "French" else "ES"
    r = requests.post(
        "https://api-free.deepl.com/v2/translate",
        data={"auth_key":DEEPL_KEY,
              "text": text,
              "target_lang": target,
              "formality": "more"}
    )
    return r.json()["translations"][0]["text"]

# ─── Parse GPT output into dict ─────────────────────────────────────────────────
def parse_sections(out: str) -> dict:
    import re

    # map any seen heading (English or French) → our internal English key
    label_map = {
        # English
        "Title":                   "Title",
        "Aecon Business Sector":   "Aecon Business Sector",
        "Project/Location":        "Project/Location",
        "Date of Event":           "Date of Event",
        "Event Type":              "Event Type",
        "Event Summary Header":    "Event Summary Header",
        "Event Summary":           "Event Summary",
        "Contributing Factors":    "Contributing Factors",
        "Lessons Learned":         "Lessons Learned",

        # French
        "Titre":                   "Title",
        "Secteur d'activité d'Aecon":    "Aecon Business Sector",
        "Secteur d’activité d’Aecon":    "Aecon Business Sector",
        "Secteur d'activité Aecon":      "Aecon Business Sector",
        "Projet/Emplacement":     "Project/Location",
        "Projet/Lieu":            "Project/Location",
        "Date de l'événement":    "Date of Event",
        "Type d'événement":       "Event Type",
        "En-tête du résumé de l'événement": "Event Summary Header",
        "En‑tête du résumé de l’événement": "Event Summary Header",
        "Résumé de l'événement":  "Event Summary",
        "Facteurs contributifs":  "Contributing Factors",
        "Leçons apprises":        "Lessons Learned",
        "Leçons tirées":          "Lessons Learned",
    }

    sections = {}
    current = None
    # capture “Label: optional value”
    pat = re.compile(r'^([^:]+):\s*(.*)$')
    for line in out.splitlines():
        line = line.strip()
        if not line:
            continue

        m = pat.match(line)
        if m:
            raw_label = m.group(1).strip()        # strip off any trailing spaces
            rest      = m.group(2).strip()
            key = label_map.get(raw_label)
            if key:
                sections[key] = []
                if rest:
                    sections[key].append(rest)
                current = key
                continue

        # if we’re inside a section, append any non‐header lines
        if current:
            sections[current].append(line)

    return sections



# ─── Render + insert images into template ──────────────────────────────────────
def render_with_docxtpl(secs: dict, tpl_path: str, out_path: str, images: list[str]):
    # verify images via PIL
    valid_exts = {'.png','.jpg','.jpeg','.bmp','.gif'}
    verified = []
    for img in images:
        ext = os.path.splitext(img)[1].lower()
        if ext not in valid_exts:
            continue
        try:
            with PILImage.open(img) as im:
                im.verify()
            verified.append(img)
        except:
            continue
    images = verified

        # filter out banner-like headers (very wide & short)
    filtered = []
    for img in images:
        try:
            with PILImage.open(img) as im:
                w, h = im.size
            # skip if width is more than 3x height (likely slide banner)
            if h == 0 or w / h > 3:
                continue
            filtered.append(img)
        except:
            continue
    images = filtered

    # Fill template
    tpl = DocxTemplate(tpl_path)
    context = {
        "TITLE":   secs.get("Title", [""])[0],
        "SECTOR":  " ".join(secs.get("Aecon Business Sector", [])),
        "PROJECT":" ".join(secs.get("Project/Location", [])),
        "DATE":   " ".join(secs.get("Date of Event", [])),
        "EVENT_TYPE":" ".join(secs.get("Event Type", [])),
        "SUMMARY_HEADER": secs.get("Event Summary Header", [""])[0],
        "SUMMARY":  " \n".join(secs.get("Event Summary", [])),
        "FACTORS":  "\n".join(f"{f}" for f in secs.get("Contributing Factors", [])),
        "LESSONS":  "\n".join(f"{l}" for l in secs.get("Lessons Learned", [])),
    }


    # Insert images
    MAX_IMG = 10
    for i in range(1, MAX_IMG+1):
        key = f"IMAGE{i}"
        if i <= len(images):
            try:
                with PILImage.open(images[i-1]) as im:
                    im.verify()
                context[key] = InlineImage(tpl, images[i-1], width=Inches(2.5))
            except:
                context[key] = ""
        else:
            context[key] = ""

    tpl.render(context)
    tpl.save(out_path)

# ─── Streamlit UI ────────────────────────────────────────────────────────────
st.set_page_config(page_title="Aecon Lessons Learned Generator", page_icon="📘")
st.image("AECON.png", width=300)
st.markdown("""
<style>
  .stApp { background:#fff; }
  h1,h2 { color:#c8102e; }
  .stButton>button,.stDownloadButton>button { background:#c8102e;color:#fff; }
  body { font-family:'Segoe UI',sans-serif; }
</style>
""", unsafe_allow_html=True)

st.title("Serious Event Lessons Learned Generator")
pptx_file = st.file_uploader("Upload Executive Review PPTX", type="pptx")
lang = st.selectbox("Language:", ["English","French (Canadian)","Spanish - waiting on template"])
translator = None
if lang in ["French (Canadian)", "Spanish"]:
    translator = st.radio("Translate via:", ["OpenAI","DeepL"])

if pptx_file and st.button("📄 Generate DOCX"):
    # 0) prepare
    os.makedirs("input", exist_ok=True)
    in_fp = f"input/{pptx_file.name}"
    with open(in_fp, "wb") as f:
        f.write(pptx_file.getbuffer())

    # initialize progress
    progress = st.progress(0)

    # 1) extract text & images
    progress.progress(10)
    raw_text, images = extract_text_and_images_from_pptx(in_fp)

    # 2) summarize
    progress.progress(30)
    generated = summarize_and_extract(raw_text)

    # show raw for debug
    st.text_area("📝 Raw Generated Content", generated, height=300)

    # 3) translate if needed
    progress.progress(50)
    if lang == "French (Canadian)":
        generated = (
            translate_to_french_openai(generated)
            if translator == "OpenAI"
            else translate_to_deepl(generated, "French")
        )
    elif lang == "Spanish":
        generated = (
            translate_to_spanish_openai(generated)
            if translator == "OpenAI"
            else translate_to_deepl(generated, "Spanish")
        )

    # show translated
    if lang != "English":
        st.text_area(f"📝 Translated ({lang})", generated, height=300)

    # 4) parse & render
    progress.progress(75)
    secs = parse_sections(generated)
    out_fp = f"lessons_learned_{lang[:2].lower()}.docx"
    if lang.startswith("French"):
        tpl_path = "lessons learned template fr.docx"
    else:
        tpl_path = "lessons learned template.docx"
    render_with_docxtpl(secs, tpl_path, out_fp, images)


    # render_with_docxtpl(secs, "lessons learned template.docx", out_fp, images)

    # 5) done
    progress.progress(100)
    st.success("✅ Done generating!")

    with open(out_fp, "rb") as f:
        st.download_button(
            "📥 Download DOCX",
            f,
            out_fp,
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
st.markdown("""
<hr style="border:none;height:2px;background:#c8102e;"/>
<div style="text-align:center;padding:10px;background:#c8102e;color:#fff;">
  Built by Aecon | For internal use only
</div>
""", unsafe_allow_html=True)
