import os
import re
import streamlit as st
import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from openai import OpenAI
from docxtpl import DocxTemplate, InlineImage
from docx import Document
from docx.shared import Inches

# ─── Load secrets ─────────────────────────────────────────────────────────────
load_dotenv()
OPENAI_KEY = os.getenv("OPENAI_API_KEY")
DEEPL_KEY  = os.getenv("DEEPL_API_KEY")
if not OPENAI_KEY:
    st.error("OpenAI API key missing. Set OPENAI_API_KEY in .env or Streamlit secrets.")
client = OpenAI(api_key=OPENAI_KEY)

# ─── Extract PPTX text & images ───────────────────────────────────────────────
def extract_text_and_images_from_pptx(path: str):
    prs = Presentation(path)
    text, images = "", []
    os.makedirs("images", exist_ok=True)

    for idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if shp.has_text_frame:
                text += shp.text_frame.text + "\n"
            elif shp.has_table:
                for row in shp.table.rows:
                    for cell in row.cells:
                        text += cell.text + "\n"
            elif shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                blob = shp.image.blob
                ext  = shp.image.ext
                fn   = f"images/slide_{idx+1}_{len(images)}.{ext}"
                with open(fn, "wb") as imgf:
                    imgf.write(blob)
                images.append(fn)

    return text, images

# ─── GPT Summarization ─────────────────────────────────────────────────────────
def summarize_and_extract(text: str) -> str:
    prompt = f"""
You are helping prepare a standardized Lessons Learned document from a serious incident.

Please extract and clearly label the following sections. Each section label should appear on its own line and be followed by its content. Separate sections with one blank line.

Use these exact labels:
Title:
Aecon Business Sector:
Project/Location:
Date of Event:
Event Type:
Event Summary:
Contributing Factors:
Lessons Learned:

Here is the presentation text:
{text}
"""
    res = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role":"user","content":prompt}],
        temperature=0.2,
        max_tokens=1000,
    )
    return res.choices[0].message.content.strip()

# ─── Translations ─────────────────────────────────────────────────────────────
def translate_to_french_openai(text: str) -> str:
    prompt = f"Translate the following into professional French Canadian, keep formatting:\n\n{text}"
    res = client.chat.completions.create(
        model="gpt-4",
        messages=[{"role":"user","content":prompt}],
        temperature=0.2,
        max_tokens=1000,
    )
    return res.choices[0].message.content.strip()

def translate_to_french_deepl(text: str) -> str:
    if not DEEPL_KEY:
        st.error("DeepL API key missing.")
        return text
    r = requests.post(
        "https://api-free.deepl.com/v2/translate",
        data={"auth_key":DEEPL_KEY,"text":text,"target_lang":"FR","formality":"more"}
    )
    try:
        return r.json()["translations"][0]["text"]
    except:
        st.error("DeepL failed")
        return text

# ─── Parse GPT output ──────────────────────────────────────────────────────────
def parse_sections(content: str) -> dict:
    sections, current = {}, None
    for line in content.splitlines():
        line = line.strip()
        if not line: continue
        if line.endswith(':') and len(line.split())<6:
            current = line[:-1].strip()
            sections[current] = []
        elif current:
            sections[current].append(line)
    return sections

# ─── Template rendering + image append ────────────────────────────────────────
def render_with_docxtpl(sections: dict, tpl_path: str, out_path: str, images: list[str]):
    tpl = DocxTemplate(tpl_path)
    context = {
        "TITLE":      sections.get("Title", [""])[0],
        "SECTOR":     " ".join(sections.get("Aecon Business Sector", [])),
        "PROJECT":    " ".join(sections.get("Project/Location", [])),
        "DATE":       " ".join(sections.get("Date of Event", [])),
        "EVENT_TYPE": " ".join(sections.get("Event Type", [])),
        "SUMMARY":    " ".join(sections.get("Event Summary", [])),
        "FACTORS":    sections.get("Contributing Factors", []),
        "LESSONS":    sections.get("Lessons Learned", []),
    }
    tpl.render(context)
    tpl.save(out_path)

    # now append images via python-docx
    doc = Document(out_path)
    if images:
        doc.add_page_break()
        doc.add_heading("Supporting Pictures", level=2)
        table = doc.add_table(rows=0, cols=2)
        for i in range(0, len(images), 2):
            row = table.add_row().cells
            for j in (0,1):
                if i+j < len(images):
                    try:
                        row[j].paragraphs[0].add_run().add_picture(
                            images[i+j], width=Inches(2.5)
                        )
                    except:
                        row[j].text = "[Image failed]"
    doc.save(out_path)

# ─── Fallback manual DOCX builds ──────────────────────────────────────────────
def create_lessons_learned_doc(content: str, out_path: str, images=None):
    secs = parse_sections(content)
    doc = Document()
    doc.add_heading(secs.get("Title",[""])[0], level=1)
    doc.add_heading("Aecon Business Sector", level=2)
    doc.add_paragraph(" ".join(secs.get("Aecon Business Sector",[])))
    doc.add_heading("Project/Location", level=2)
    doc.add_paragraph(" ".join(secs.get("Project/Location",[])))
    doc.add_heading("Date of Event", level=2)
    doc.add_paragraph(" ".join(secs.get("Date of Event",[])))
    doc.add_heading("Event Type", level=2)
    doc.add_paragraph(" ".join(secs.get("Event Type",[])))
    doc.add_heading("Event Summary", level=2)
    doc.add_paragraph(" ".join(secs.get("Event Summary",[])))
    doc.add_heading("Contributing Factors", level=2)
    for f in secs.get("Contributing Factors",[]): doc.add_paragraph(f,style="List Bullet")
    doc.add_heading("Lessons Learned", level=2)
    for l in secs.get("Lessons Learned",[]): doc.add_paragraph(l,style="List Bullet")
    if images:
        doc.add_page_break()
        doc.add_heading("Supporting Pictures", level=2)
        tbl = doc.add_table(rows=0, cols=2)
        for i in range(0,len(images),2):
            row = tbl.add_row().cells
            for j in (0,1):
                if i+j < len(images):
                    try:
                        row[j].paragraphs[0].add_run().add_picture(
                            images[i+j], width=Inches(2.5)
                        )
                    except:
                        row[j].text = "⚠️ Error"
    doc.save(out_path)

def create_lessons_learned_doc_fr(content: str, out_path: str, images=None):
    secs = parse_sections(content)
    doc = Document()
    G = secs.get
    title = G("Titre",[""])[0]
    doc.add_heading(title or "Document d'apprentissage", level=1)
    doc.add_heading("Secteur d'activité d'Aecon", level=2)
    doc.add_paragraph(" ".join(G("Secteur d'activité d'Aecon",[])))
    doc.add_heading("Projet/Emplacement", level=2)
    doc.add_paragraph(" ".join(G("Projet/Emplacement",[])))
    doc.add_heading("Date de l'événement", level=2)
    doc.add_paragraph(" ".join(G("Date de l'événement",[])))
    doc.add_heading("Type d'événement", level=2)
    doc.add_paragraph(" ".join(G("Type d'événement",[])))
    doc.add_heading("Résumé de l'événement", level=2)
    doc.add_paragraph(" ".join(G("Résumé de l'événement",[])))
    doc.add_heading("Facteurs contributifs", level=2)
    for f in G("Facteurs contributifs",[]): doc.add_paragraph(f,style="List Bullet")
    doc.add_heading("Leçons apprises", level=2)
    for l in G("Leçons apprises",[])+G("Leçons tirées",[]): doc.add_paragraph(l,style="List Bullet")
    if images:
        doc.add_page_break()
        doc.add_heading("Photos de soutien", level=2)
        tbl = doc.add_table(rows=0, cols=2)
        for i in range(0,len(images),2):
            row = tbl.add_row().cells
            for j in (0,1):
                if i+j < len(images):
                    try:
                        row[j].paragraphs[0].add_run().add_picture(
                            images[i+j], width=Inches(2.5)
                        )
                    except:
                        row[j].text = "⚠️ Erreur"
    doc.save(out_path)

# ─── Streamlit UI ────────────────────────────────────────────────────────────
st.set_page_config(page_title="Aecon Lessons Learned Generator", page_icon="📘")
st.image("AECON.png", width=300)
st.markdown("""
<style>
  .stApp { background:#fff; }
  h1,h2 { color:#c8102e; }
  .stButton>button, .stDownloadButton>button { background:#c8102e; color:#fff; }
  body { font-family:'Segoe UI',sans-serif; }
</style>
""", unsafe_allow_html=True)

st.title("🦺 Serious Event Lessons Learned Generator")
uploaded = st.file_uploader("Upload Executive Review PPTX", type="pptx")
language = st.selectbox("Language:", ["English","French (Canadian)"])
translator = None
if language=="French (Canadian)":
    translator = st.radio("Translate via:", ["OpenAI","DeepL"])
use_tpl = st.checkbox("📑 Use official Word template", True)

if uploaded and st.button("📄 Generate Document"):
    # Save PPTX
    in_fp = f"input/{uploaded.name}"
    os.makedirs("input", exist_ok=True)
    with open(in_fp,"wb") as f: f.write(uploaded.getbuffer())

    # Extract + Debug raw text
    raw_text, images = extract_text_and_images_from_pptx(in_fp)
    st.subheader("🔍 Raw PPTX Text")
    st.text_area("", raw_text, height=200)

    # Summarize + Debug GPT output
    with st.spinner("Summarizing..."):
        generated = summarize_and_extract(raw_text)
    st.subheader("✍️ GPT Output")
    st.text_area("", generated, height=300)

    # Translate if needed
    if language=="French (Canadian)":
        with st.spinner("Translating..."):
            generated = (translate_to_french_openai(generated)
                         if translator=="OpenAI"
                         else translate_to_french_deepl(generated))

    # Parse & render
    secs = parse_sections(generated)
    out_fp = f"lessons_learned_{'fr' if language.startswith('French') else 'en'}.docx"

    if use_tpl:
        render_with_docxtpl(secs, "template.docx", out_fp, images)
    else:
        if language.startswith("French"):
            create_lessons_learned_doc_fr(generated, out_fp, images)
        else:
            create_lessons_learned_doc(generated, out_fp, images)

    # Download
    with open(out_fp,"rb") as f:
        st.download_button(
            "📥 Download DOCX", f, out_fp,
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

st.markdown("""
<hr style="border:none;height:2px;background:#c8102e;"/>
<div style="text-align:center;padding:10px;background:#c8102e;color:#fff;">
  Built by Aecon | For internal use only
</div>
""", unsafe_allow_html=True)
