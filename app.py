import os
import streamlit as st
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Pt, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from dotenv import load_dotenv
from openai import OpenAI
import re

# Load environment variables
load_dotenv()

# Set up OpenAI client
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

# Extract text and images from PPTX

def extract_text_and_images_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ''
    image_paths = []
    os.makedirs("images", exist_ok=True)

    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:
                text += shape.text_frame.text + "\n"
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                image_bytes = image.blob
                image_path = f"images/slide_{i+1}_{len(image_paths)}.{ext}"
                with open(image_path, 'wb') as img_file:
                    img_file.write(image_bytes)
                image_paths.append(image_path)

    return text, image_paths

# Generate summary and extracted sections using OpenAI GPT

def summarize_and_extract(text):
    prompt = f"""
You are helping prepare a standardized Lessons Learned document from a serious incident.

Please extract and clearly label the following sections. Make sure each section label is written on its own line, followed by its content, with a blank line between sections.

Use these exact labels:
Title:
Aecon Business Sector:
Project/Location:
Date of Event:
Event Type:
Event Summary:
Contributing Factors:
Lessons Learned:

Here is the presentation text:
{text}

Output format:
Title:
[...title here...]

Aecon Business Sector:
[...]

Project/Location:
[...]

...and so on.
"""


    response = client.chat.completions.create(
        model="gpt-3.5-turbo",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.2,
        max_tokens=1000,
    )

    return response.choices[0].message.content.strip()

# Generate Word Document from the summarized content (no template needed)


def create_lessons_learned_doc(content, output_path, image_paths=None):
    doc = Document()

    # Helper function to match section keys
    def get_section_by_keyword(sections, keyword, default=''):
        for key in sections:
            if keyword.lower() in key.lower():
                return sections[key]
        return [default] if default else []

        
    
    

    def parse_sections(content):
        sections = {}
        current_section = None
        for line in content.splitlines():
            line = line.strip()
            if not line:
                continue
            if line.endswith(':') and len(line.split()) < 6:  # likely a section header
                current_section = line[:-1]
                sections[current_section] = []
            elif current_section:
                sections[current_section].append(line)
        return sections

    sections = parse_sections(content)

    # Title
    title = get_section_by_keyword(sections, "title", "Lessons Learned Handout")[0]
    doc.add_heading(title, level=1)

    # Meta Sections
    doc.add_heading("Aecon Business Sector", level=2)
    sector = ' '.join(sections.get("Aecon Business Sector",[]))
    doc.add_paragraph(sector)

    doc.add_heading("Project/Location", level=2)
    location = ' '.join(sections.get("Project/Location",[]))
    doc.add_paragraph(location)

    doc.add_heading("Date of Event", level=2)
    date = ' '.join(sections.get("Date of Event",[]))
    doc.add_paragraph(date)

    doc.add_heading("Event Type", level=2)
    event = ' '.join(sections.get("Event Type", []))
    doc.add_paragraph(event)

    # Event Summary
    doc.add_heading("Event Summary", level=2)
    summary = ' '.join(sections.get("Event Summary", []))
    doc.add_paragraph(summary)

    # Contributing Factors
    doc.add_heading("Contributing Factors", level=2)
    for factor in get_section_by_keyword(sections, "contributing"):
        doc.add_paragraph(factor, style='List Bullet')

    # Lessons Learned
    doc.add_heading("Lessons Learned", level=2)
    for lesson in get_section_by_keyword(sections, "lessons"):
        doc.add_paragraph(lesson, style='List Bullet')

    # Supporting Pictures
    if image_paths:
        doc.add_page_break()
        doc.add_heading("Supporting Pictures", level=2)

        table = doc.add_table(rows=0, cols=2)
        table.autofit = True

        for i in range(0, len(image_paths), 2):
            row_cells = table.add_row().cells
            for j in range(2):
                if i + j < len(image_paths):
                    try:
                        paragraph = row_cells[j].paragraphs[0]
                        run = paragraph.add_run()
                        run.add_picture(image_paths[i + j], width=Inches(2.5))
                    except Exception:
                        row_cells[j].text = "⚠️ Error loading image"

    # Save the final document
    doc.save(output_path)


# Streamlit UI
st.title('🦺 Serious Event Lessons Learned Generator')

uploaded_file = st.file_uploader("Upload Executive Review PPTX", type=['pptx'])

if uploaded_file:
    input_filepath = f'input/{uploaded_file.name}'
    os.makedirs("input", exist_ok=True)
    with open(input_filepath, 'wb') as f:
        f.write(uploaded_file.getbuffer())

    with st.spinner("Extracting content from presentation..."):
        pptx_text, extracted_images = extract_text_and_images_from_pptx(input_filepath)

    with st.spinner("Generating Lessons Learned summary..."):
        generated_content = summarize_and_extract(pptx_text)

    st.success("✅ Generation Complete!")

    st.text_area("📝 Generated Content:", generated_content, height=300)

    output_path = 'generated_lessons_learned.docx'
    create_lessons_learned_doc(generated_content, output_path, extracted_images)

    with open(output_path, "rb") as file:
        st.download_button(
            label="📥 Download Lessons Learned DOCX",
            data=file,
            file_name="Lessons_Learned_Summary.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )